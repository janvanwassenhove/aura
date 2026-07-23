"""U205: build the test presentation deck to accompany the scenario.

Run:  uv run --with python-pptx python docs/demo/build_robot_junior_dev_pptx.py

Slide numbers here line up with the `slide:N` triggers in
robot-junior-dev.scenario.yaml (PowerPoint numbers slides from 1).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

SLIDES = [
    ("I Hired a Real Robot as My Junior Dev",
     "…and now my kids don't want to learn Java anymore.\n\n"
     "A stand-up, a field report, and one uncomfortable question."),
    ("For decades, we taught the hard part",
     "Java, C#, Python, patterns, frameworks.\n"
     "Because writing software WAS the hard part."),
    ("My son isn't learning Java",
     "He's learning to make money.\n"
     "A study app one week, a gaming overlay the next.\n"
     "He ships from ideas, not code. He sees products, not 'programming'."),
    ("So I hired a robot",
     "Reachy Mini — a real open-source desktop robot — as my junior dev.\n"
     "I direct it; it drives a fleet of agents that code while I sleep.\n\n"
     "My bet: writing software becomes a commodity. Expertise doesn't."),
    ("The interface is shifting",
     "From typing to talking.\n"
     "Directing agents in plain language — until speaking to a robot\n"
     "on your desk beats typing into an editor."),
    ("Are we raising the last generation of developers?",
     "…or the first generation of something else?"),
    ("What still matters",
     "Architecture. Systems thinking. Integration. Taste.\n"
     "Knowing WHAT to build.\n\n"
     "There will be live robot moments. Some may even work."),
]


def build(out: Path) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content
    for title, body in SLIDES:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        for i, line in enumerate(body.split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
            para.font.size = Pt(20)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"wrote {out} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    build(Path(__file__).with_name("robot-junior-dev.pptx"))
